#!/usr/bin/env python3
"""Generate Credara pitch decks (PPTX): judge edition + investor edition."""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Credara_Pitch_Deck.pptx"
OUT_INVESTOR = ROOT / "docs" / "Credara_Investor_Deck.pptx"
LANDING_SHOT = ROOT / "docs" / "screenshots" / "landing.png"
WORKSPACE_SHOT = ROOT / "docs" / "screenshots" / "workspace-judge-mode.png"
POLYGONSCAN_SHOT = ROOT / "docs" / "screenshots" / "polygonscan-tx.png"
RELEASE_TX_URL = "https://amoy.polygonscan.com/tx/0x5743a885e54063da6c4056d73e4b49113918558fd09d8973c9074de8e57af9de"
PROOF_TX_URL = "https://amoy.polygonscan.com/tx/0x71b1d6c74b30033b7f3ab1cad174bf75cf2003fce593dd8c7e04aa3964acf251"

# Brand palette
NAVY = RGBColor(0x07, 0x14, 0x2F)
INK = RGBColor(0x11, 0x18, 0x27)
BLUE = RGBColor(0x25, 0x63, 0xEB)
MUTED = RGBColor(0x66, 0x70, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x16, 0xA3, 0x4A)


def _set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _add_bullets(slide, left, top, width, height, items, *, size=16, color=INK, spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(size * spacing * 0.35)
    return box


def _add_footer(slide, *, dark: bool = False) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY if dark else RGBColor(0xF1, 0xF5, 0xF9)
    bar.line.fill.background()
    text_color = WHITE if dark else MUTED
    _add_textbox(slide, Inches(0.7), Inches(7.12), Inches(6), Inches(0.3), "Credara · Smart Commerce Infrastructure", size=10, color=text_color)
    _add_textbox(slide, Inches(9.5), Inches(7.12), Inches(3.5), Inches(0.3), "credara-jet.vercel.app", size=10, color=text_color, align=PP_ALIGN.RIGHT)


def _accent_line(slide, top: float = 1.75) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(top), Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def _section_header(slide, title: str, subtitle: str = "") -> None:
    _set_slide_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.7), Inches(0.45), Inches(11.5), Inches(0.5), title, size=11, bold=True, color=BLUE)
    _add_textbox(slide, Inches(0.7), Inches(0.85), Inches(11.5), Inches(1.2), subtitle or title, size=32, bold=True, color=INK)
    _accent_line(slide)
    _add_footer(slide)


def _add_stat_card(slide, left, top, width, height, value: str, label: str) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = SOFT
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.55), value, size=22, bold=True, color=BLUE)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.62), width - Inches(0.3), Inches(0.45), label, size=11, color=MUTED)


def _add_comparison_row(slide, y: float, name: str, pain: str, credara: str, *, highlight: bool = False) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(0.72))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF) if highlight else WHITE
    bg.line.color.rgb = RGBColor(0xE7, 0xEA, 0xF1)
    name_color = BLUE if highlight else INK
    _add_textbox(slide, Inches(0.85), Inches(y + 0.14), Inches(2.3), Inches(0.45), name, size=13, bold=True, color=name_color)
    _add_textbox(slide, Inches(3.2), Inches(y + 0.1), Inches(4.5), Inches(0.55), pain, size=12, color=MUTED)
    _add_textbox(slide, Inches(7.9), Inches(y + 0.1), Inches(4.5), Inches(0.55), credara, size=12, bold=highlight, color=INK if highlight else MUTED)


def _add_aha_climax_slide(slide, *, investor: bool = False) -> None:
    """Demo climax — money moving on Polygon, not slideware."""
    _section_header(slide, "THE AHA MOMENT", "Supplier paid after verified delivery — not after 45 days of paperwork")
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(2.15),
        Inches(5.8),
        Inches(1.35),
        "1,000 MockUSDC released on Polygon Amoy after Smart LC verified delivery.\nBank LC would still be on day 3 of manual processing.",
        size=18,
        bold=True,
        color=INK,
    )
    if investor:
        _add_textbox(
            slide,
            Inches(0.9),
            Inches(3.55),
            Inches(5.8),
            Inches(0.9),
            "Example: AED 250k invoice → 80% advance → funded same afternoon → released after delivery verify.",
            size=15,
            color=MUTED,
        )
    else:
        _add_bullets(
            slide,
            Inches(0.9),
            Inches(3.55),
            Inches(5.8),
            Inches(2.2),
            [
                "Live product: credara-jet.vercel.app/workspace",
                "Judge Mode walks the flow in ~3 minutes",
                "Fail-closed: Simulated vs Anchored — no fake links",
            ],
            size=14,
        )
    if POLYGONSCAN_SHOT.exists():
        slide.shapes.add_picture(str(POLYGONSCAN_SHOT), Inches(7.0), Inches(2.1), width=Inches(5.8))
    climax = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.75))
    climax.fill.solid()
    climax.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    climax.line.color.rgb = GREEN
    _add_textbox(
        slide,
        Inches(1.05),
        Inches(6.0),
        Inches(11.2),
        Inches(0.5),
        f"Verified release tx: {RELEASE_TX_URL}",
        size=11,
        bold=True,
        color=GREEN,
    )


def _add_before_after_slide(slide) -> None:
    _section_header(slide, "BEFORE · AFTER", "Collapse 7–10 day LC friction → same-day proof-to-funding")
    _add_textbox(slide, Inches(0.9), Inches(2.05), Inches(5.5), Inches(0.4), "Traditional bank LC", size=14, bold=True, color=MUTED)
    bank = [
        "7–10 days manual processing",
        "Paper trails banks can't verify fast",
        "~40% SME applications rejected",
        "Supplier waits 45–90 days to get paid",
    ]
    y = 2.5
    for line in bank:
        _add_textbox(slide, Inches(0.9), Inches(y), Inches(5.5), Inches(0.45), f"✗  {line}", size=15, color=MUTED)
        y += 0.62
    _add_textbox(slide, Inches(7.0), Inches(2.05), Inches(5.5), Inches(0.4), "Credara (live today)", size=14, bold=True, color=BLUE)
    credara = [
        "Buyer confirms + delivery proof verified",
        "Proof hash anchored on Polygon",
        "Financier funds Smart LC escrow",
        "Stablecoin released same day after verify",
    ]
    y = 2.5
    for line in credara:
        _add_textbox(slide, Inches(7.0), Inches(y), Inches(5.5), Inches(0.45), f"✓  {line}", size=15, bold=True, color=INK)
        y += 0.62
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.75), Inches(1.95), Inches(0.05), Inches(4.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE7, 0xEA, 0xF1)
    rect.line.fill.background()
    _add_stat_card(slide, Inches(0.9), Inches(5.35), Inches(2.6), Inches(1.05), "7–10 days", "Bank LC")
    _add_stat_card(slide, Inches(3.65), Inches(5.35), Inches(2.6), Inches(1.05), "45–90 days", "SME wait")
    _add_stat_card(slide, Inches(6.4), Inches(5.35), Inches(2.6), Inches(1.05), "Same day", "Proof-to-fund")
    _add_stat_card(slide, Inches(9.15), Inches(5.35), Inches(2.6), Inches(1.05), "3 min", "Judge Mode demo")


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Title
    slide = prs.slides.add_slide(blank)
    _set_slide_bg(slide, NAVY)
    _add_textbox(slide, Inches(0.9), Inches(1.45), Inches(11), Inches(1.0), "Credara", size=54, bold=True, color=WHITE)
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(2.55),
        Inches(11),
        Inches(1.0),
        "Problem #1: SME trade finance on Polygon —\ntokenized receivables, Smart LC, on-chain credit scoring",
        size=22,
        bold=True,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(3.75),
        Inches(11),
        Inches(0.9),
        "Smart Commerce Infrastructure Challenge  ·  Powered by Polygon  ·  DIFC-ready UAE corridor",
        size=14,
        color=MUTED,
    )
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(4.65),
        Inches(11),
        Inches(0.55),
        "UAE: $50B+ payments corridor · PTSR live · 51.3% stablecoin activity — B2B settlement wedge",
        size=12,
        color=RGBColor(0x94, 0xA3, 0xB8),
    )
    _add_textbox(slide, Inches(0.9), Inches(5.85), Inches(11), Inches(0.5), "Othniel Obasi & Martins Nwanu  ·  NOVTIA Ltd  ·  credara-jet.vercel.app", size=14, color=RGBColor(0x94, 0xA3, 0xB8))
    _add_footer(slide, dark=True)

    # 2 — Aha climax (open here when pitching live)
    slide = prs.slides.add_slide(blank)
    _add_aha_climax_slide(slide, investor=False)

    # 3 — Before / After
    slide = prs.slides.add_slide(blank)
    _add_before_after_slide(slide)

    # 4 — Challenge fit (maps to official brief)
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "CHALLENGE FIT", "Built exactly what Problem #1 asks for — on Polygon, live today")
    _add_textbox(slide, Inches(0.9), Inches(2.05), Inches(5.8), Inches(0.4), "Hackathon requirement", size=12, bold=True, color=BLUE)
    _add_textbox(slide, Inches(6.9), Inches(2.05), Inches(5.8), Inches(0.4), "Credara (shipped)", size=12, bold=True, color=GREEN)
    pairs = [
        ("Tokenized receivables", "ReceivableRegistry + buyer-confirmed invoices → finance-ready receivables"),
        ("Smart contract letters of credit", "SmartLCFactory / SmartLC escrow lifecycle on Amoy"),
        ("On-chain trade credit scoring", "CreditScoreAttestation + scoring workers"),
        ("Unlock SME capital · reduce friction", "Judge Mode: proof anchor → fund LC → release in stablecoin"),
    ]
    y = 2.5
    for req, delivered in pairs:
        _add_textbox(slide, Inches(0.9), Inches(y), Inches(5.6), Inches(0.55), f"✓  {req}", size=15, bold=True, color=INK)
        _add_textbox(slide, Inches(6.9), Inches(y), Inches(5.8), Inches(0.55), delivered, size=14, color=MUTED)
        y += 0.78
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.75))
    note.fill.solid()
    note.fill.fore_color.rgb = SOFT
    note.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    _add_textbox(
        slide,
        Inches(1.05),
        Inches(5.9),
        Inches(11.2),
        Inches(0.5),
        "Problem #2 (POS / tourist wallets) is out of MVP scope — same settlement rails are PTSR-ready for licensed AED stablecoin B2B use.",
        size=13,
        color=INK,
    )

    # 3 — Judging criteria
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "JUDGING CRITERIA", "How Credara maps to what the panel scores")
    criteria = [
        ("Innovation", "Smart LC + proof anchoring + receivables in one UAE corridor stack"),
        ("Technical excellence", "Slither CI, fail-closed chain, live deploy, role-separated contracts"),
        ("User experience", "4 enterprise workspaces + 3-min Judge Mode guided demo"),
        ("Real-world impact", "$2T gap · 40% SME rejection · Jebel Ali 15M+ TEUs — suppliers paid faster"),
    ]
    positions = [(0.9, 2.15), (6.9, 2.15), (0.9, 4.35), (6.9, 4.35)]
    for (left, top), (title, body) in zip(positions, criteria):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.6), Inches(1.85))
        card.fill.solid()
        card.fill.fore_color.rgb = SOFT
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        _add_textbox(slide, Inches(left + 0.15), Inches(top + 0.15), Inches(5.3), Inches(0.4), title, size=15, bold=True, color=BLUE)
        _add_textbox(slide, Inches(left + 0.15), Inches(top + 0.55), Inches(5.3), Inches(1.1), body, size=13, color=INK)

    # 4 — Problem
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "PROBLEM", "SMEs wait months to get paid on goods already delivered")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(11.5),
        Inches(4.5),
        [
            "$2T+ global trade finance gap — demand unfilled or priced out",
            "~40% of SME trade-finance applications rejected (weak evidence, not weak fundamentals)",
            "Traditional LCs: 7–10 days of manual bank coordination and paper",
            "Invoice fraud & duplicate financing — lenders can't verify delivery independently",
            "UAE corridor: Jebel Ali 15M+ TEUs/yr — high volume, high friction for suppliers",
        ],
        size=18,
    )

    # 5 — Solution
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "SOLUTION", "Verified trade evidence → finance-ready receivables → Smart LC settlement")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(11.5),
        Inches(4.5),
        [
            "Buyer-confirmed invoices + delivery proof become finance-ready receivables",
            "Proof hashes anchored on Polygon — tamper-evident, documents stay off-chain",
            "Programmable Smart LC escrow: fund → verify delivery → release in stablecoin",
            "Four workspaces (SME, Buyer, Financier, Admin) on one shared trust layer",
            "Wedge: SME trade-finance slice of UAE corridor — not consumer remittance or POS",
        ],
        size=18,
    )

    # 6 — How it works
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "HOW IT WORKS", "End-to-end flow (~3 min Judge Mode)")
    steps = [
        ("1", "SME", "Order → invoice → delivery proof"),
        ("2", "Buyer", "Confirm trade obligation"),
        ("3", "SME", "Create receivable · anchor proof on Polygon"),
        ("4", "Financier", "Deal room review · fund Smart LC"),
        ("5", "System", "Verify delivery · release settlement to seller"),
    ]
    y = 2.1
    for num, actor, action in steps:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y), Inches(0.5), Inches(0.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BLUE
        rect.line.fill.background()
        _add_textbox(slide, Inches(0.9), Inches(y + 0.06), Inches(0.5), Inches(0.4), num, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(slide, Inches(1.55), Inches(y + 0.05), Inches(2.0), Inches(0.35), actor, size=13, bold=True, color=BLUE)
        _add_textbox(slide, Inches(3.6), Inches(y + 0.02), Inches(8.8), Inches(0.55), action, size=16, color=INK)
        y += 0.92
    _add_textbox(slide, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4), "Live: credara-jet.vercel.app/workspace", size=12, color=MUTED)

    # 7 — Differentiators
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "WHY CREDARA", "Institutional-grade infrastructure, not a hash-on-chain demo")
    cols = [
        ("Fail-closed chain", "No fake Polygonscan links. Simulated only when tx isn't on-chain."),
        ("Real persistence", "Postgres-backed workspaces — no canned seed data for judges."),
        ("Role-separated LC", "Verifier, dispute-resolver, pauser — not one key for every escrow."),
        ("Evidence-first", "AIRG/SURGE discipline: cryptographic receipts, deterministic hashing."),
        ("Production path", "Auth0, Didit KYB, Slither CI, Vercel deploy — shipped today."),
    ]
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(11.5),
        Inches(4.8),
        [f"{t} — {b}" for t, b in cols],
        size=17,
    )

    # 8 — Architecture
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "ARCHITECTURE", "Monorepo: web · API · workers · contracts")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.0),
        Inches(5.5),
        Inches(4.5),
        [
            "apps/web — Next.js enterprise dashboard",
            "apps/api — FastAPI, Postgres, JWT, Auth0",
            "apps/workers — relayer, indexer, scoring",
            "contracts — Solidity 0.8.28, Amoy / PoS-ready",
        ],
        size=16,
    )
    _add_bullets(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.8),
        Inches(4.5),
        [
            "ProofRegistry — anchor proof hashes",
            "ReceivableRegistry — finance-ready state",
            "SmartLCFactory / SmartLC — escrow lifecycle",
            "CreditScoreAttestation — trade credit on-chain",
            "MockUSDC — demo settlement (IERC20-agnostic)",
        ],
        size=16,
        color=INK,
    )
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.85), Inches(0.05), Inches(4.8))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE7, 0xEA, 0xF1)
    rect.line.fill.background()

    # 9 — Stablecoin corridor
    slide = prs.slides.add_slide(blank)
    _set_slide_bg(slide, SOFT)
    _add_textbox(slide, Inches(0.9), Inches(0.6), Inches(11), Inches(0.4), "SETTLEMENT ROADMAP", size=11, bold=True, color=BLUE)
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(1.1),
        Inches(11),
        Inches(1.4),
        "Today: MockUSDC on Amoy.\nTomorrow: regulated AED stablecoin — same Smart LC rails.",
        size=30,
        bold=True,
        color=INK,
    )
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(3.0),
        Inches(11),
        Inches(3.5),
        [
            "Demo proves create → fund → verify → release without claiming AED issuance",
            "UAE CBUAE PTSR: regulated AED stablecoin partner — same Smart LC rails",
            "Smart LC is token-agnostic — factory and proof conditions unchanged",
            "Do say: settlement-ready · Don't say: we issue AED or we are remittance",
        ],
        size=17,
        color=INK,
    )

    _add_footer(slide)

    # 10 — Market & GTM
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "GO-TO-MARKET", "Pilot one corridor, then expand")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(5.8),
        Inches(4.5),
        [
            "Target: UAE SMEs, buyers, financiers — Jebel Ali supply chains",
            "Land financiers first — verified receivables pull SME supply",
            "DIFC ecosystem: 5,000+ VCs, 289 banks, 400+ mentors for pilot cohort",
            "Prove 7–10 day LC friction → same-day proof-to-funding",
        ],
        size=17,
    )
    _add_textbox(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.5), "Roadmap", size=14, bold=True, color=BLUE)
    _add_bullets(
        slide,
        Inches(7.2),
        Inches(2.5),
        Inches(5.5),
        Inches(4.0),
        [
            "Now: Amoy demo + Judge Mode (shipped)",
            "Near-term: Didit KYB prod, Solidity audit",
            "Phase 2: AED stablecoin partner",
            "Phase 3: Mainnet + API ecosystem",
        ],
        size=15,
        color=MUTED,
    )

    # 11 — Revenue
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "BUSINESS MODEL", "Infrastructure economics — multiple revenue lines")
    streams = [
        ("Origination fee", "Per receivable financed via Smart LC marketplace"),
        ("Platform / SaaS", "Verification, deal room, reporting — tiered by volume"),
        ("Settlement fee", "Basis points per LC fund/release vs. days of bank LC"),
        ("Credit scoring API", "Attestations for external lenders without full workspace"),
    ]
    y = 2.2
    for title, desc in streams:
        _add_textbox(slide, Inches(0.9), Inches(y), Inches(3.2), Inches(0.4), title, size=16, bold=True, color=BLUE)
        _add_textbox(slide, Inches(4.2), Inches(y), Inches(8.3), Inches(0.5), desc, size=16, color=INK)
        y += 0.85
    _add_textbox(slide, Inches(0.9), Inches(5.8), Inches(11), Inches(0.6), "API-first: third parties integrate proof, receivable, and LC rails without Credara UI.", size=14, color=MUTED)

    # 12 — Team
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "TEAM", "Banking operators + engineering leadership — built to win financier trust")
    _add_textbox(slide, Inches(0.9), Inches(2.05), Inches(5.6), Inches(0.4), "Othniel Obasi — Founder, NOVTIA Ltd (London)", size=14, bold=True, color=BLUE)
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.45),
        Inches(5.6),
        Inches(4.2),
        [
            "Digital banking: Creditville MFBank 40K → 100K+ users",
            "Utilities: PHED — 5M+ customers, $96M+ revenue growth",
            "AI governance: AIRG + SURGE v2 cryptographic receipt chain",
            "Credara: evidence-first, fail-closed on-chain settlement",
        ],
        size=15,
    )
    _add_textbox(slide, Inches(6.9), Inches(2.05), Inches(5.6), Inches(0.4), "Martins Nwanu — Head of IT, Creditville Group", size=14, bold=True, color=BLUE)
    _add_bullets(
        slide,
        Inches(6.9),
        Inches(2.45),
        Inches(5.6),
        Inches(4.2),
        [
            "Led 11-engineer team; core banking & loan platform delivery",
            "C Money app: 30K → 100K+ users; CI/CD and Agile at scale",
            "Loan disbursement volumes scaled to ₦800M–₦1B/month",
            "Credara: enterprise workspace, API, and production deploy rigor",
        ],
        size=15,
    )
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.65), Inches(1.95), Inches(0.05), Inches(4.8))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE7, 0xEA, 0xF1)
    rect.line.fill.background()

    # 13 — MVP live today
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "MVP — LIVE TODAY", "Test it now")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.0),
        Inches(5.8),
        Inches(3.8),
        [
            "credara-jet.vercel.app/workspace",
            "credara-api.vercel.app/docs",
            "Judge demo mode in sidebar",
            "Polygon Amoy contracts deployed",
            "Slither + Hardhat tests in CI",
        ],
        size=16,
    )
    if WORKSPACE_SHOT.exists():
        slide.shapes.add_picture(str(WORKSPACE_SHOT), Inches(7.0), Inches(2.0), width=Inches(5.8))
        _add_textbox(slide, Inches(7.0), Inches(6.55), Inches(5.8), Inches(0.3), "Enterprise workspace (Judge Mode)", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    else:
        _add_textbox(slide, Inches(7.5), Inches(2.0), Inches(5.2), Inches(0.5), "Non-goals (by design)", size=14, bold=True, color=BLUE)
        _add_bullets(
            slide,
            Inches(7.5),
            Inches(2.5),
            Inches(5.2),
            Inches(3.5),
            [
                "No consumer remittance MVP",
                "No POS product",
                "No AED issuance claim",
            ],
            size=16,
            color=MUTED,
        )
    _add_textbox(slide, Inches(0.9), Inches(5.5), Inches(5.8), Inches(0.5), "Non-goals: no remittance MVP · no POS · no AED issuance claim", size=12, color=MUTED)
    _add_textbox(slide, Inches(0.9), Inches(6.0), Inches(11), Inches(0.5), "Demo climax: Polygonscan proof anchor + Smart LC settlement", size=14, bold=True, color=GREEN)

    # 14 — Close
    slide = prs.slides.add_slide(blank)
    _set_slide_bg(slide, NAVY)
    _add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11), Inches(1.0), "Get paid in days, not months.", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(3.6),
        Inches(11),
        Inches(1.0),
        "Verified invoices. Polygon proof. Smart LC settlement.",
        size=20,
        color=RGBColor(0xCB, 0xD5, 0xE1),
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(slide, Inches(0.9), Inches(5.2), Inches(11), Inches(0.5), "credara-jet.vercel.app  ·  othnielobasi@gmail.com", size=14, color=MUTED, align=PP_ALIGN.CENTER)
    _add_footer(slide, dark=True)

    prs.save(OUT)
    return OUT


def build_investor() -> Path:
    """15-slide deck tuned for angels, VCs, and strategic partners."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Title (killer insight)
    slide = prs.slides.add_slide(blank)
    _set_slide_bg(slide, NAVY)
    _add_textbox(slide, Inches(0.9), Inches(1.35), Inches(11), Inches(0.8), "Credara", size=48, bold=True, color=WHITE)
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(2.25),
        Inches(11),
        Inches(1.35),
        "The LC layer banks can't ship fast enough —\nsame-day financiable receivables on Polygon.",
        size=26,
        bold=True,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(3.85),
        Inches(11),
        Inches(0.7),
        "Verified B2B trade evidence → Smart LC stablecoin settlement · UAE corridor wedge",
        size=16,
        color=MUTED,
    )
    if LANDING_SHOT.exists():
        slide.shapes.add_picture(str(LANDING_SHOT), Inches(8.6), Inches(4.55), width=Inches(3.9))
    _add_textbox(slide, Inches(0.9), Inches(5.85), Inches(7.5), Inches(0.4), "Othniel Obasi & Martins Nwanu · NOVTIA Ltd · credara-jet.vercel.app", size=13, color=RGBColor(0x94, 0xA3, 0xB8))
    _add_footer(slide, dark=True)

    # 2 — Aha climax
    slide = prs.slides.add_slide(blank)
    _add_aha_climax_slide(slide, investor=True)

    # 3 — Before / After
    slide = prs.slides.add_slide(blank)
    _add_before_after_slide(slide)

    # 4 — Problem (investor language)
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "PROBLEM", "Suppliers finance the world's trade — and wait months to get paid")
    _add_stat_card(slide, Inches(0.9), Inches(2.15), Inches(2.6), Inches(1.15), "$2T+", "Unfilled trade finance demand")
    _add_stat_card(slide, Inches(3.65), Inches(2.15), Inches(2.6), Inches(1.15), "~40%", "SME applications rejected")
    _add_stat_card(slide, Inches(6.4), Inches(2.15), Inches(2.6), Inches(1.15), "7–10 days", "Traditional LC processing")
    _add_stat_card(slide, Inches(9.15), Inches(2.15), Inches(2.6), Inches(1.15), "45–90 days", "Typical SME payment wait")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(3.55),
        Inches(11.5),
        Inches(3.2),
        [
            "Banks reject SMEs with weak paper trails — not weak businesses",
            "Invoice fraud and duplicate financing make lenders blind to real delivery",
            "UAE/Jebel Ali corridor: massive volume, same friction for suppliers",
            "Consumer remittance and POS are crowded — the SME trade-finance slice is underserved",
        ],
        size=17,
    )

    # 3 — Solution (outcomes, not stack)
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "SOLUTION", "Turn delivery proof into finance-ready receivables — settle in hours, not months")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(11.5),
        Inches(4.5),
        [
            "Buyer confirms the invoice; delivery proof is verified — receivable becomes financiable",
            "Proof anchored on Polygon (documents stay private; hashes are tamper-evident)",
            "Financier funds programmable Smart LC escrow; release only after verified delivery",
            "Four enterprise workspaces — SME, Buyer, Financier, Admin — one shared trust record",
            "Outcome: collapse 7–10 day LC friction toward same-day proof-to-funding in one corridor",
        ],
        size=18,
    )

    # 4 — Why now / Why us
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "WHY NOW · WHY US", "Regulation, market pain, and a shipped stack converge")
    _add_textbox(slide, Inches(0.9), Inches(2.05), Inches(5.5), Inches(0.4), "Why now", size=14, bold=True, color=BLUE)
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.45),
        Inches(5.8),
        Inches(4.0),
        [
            "UAE CBUAE PTSR opens regulated AED stablecoin B2B settlement",
            "Polygon enterprise infra: low-cost escrow + proof anchoring at scale",
            "SME trade-finance gap widening as corridors digitize faster than bank LC rails",
            "Smart Commerce / DIFC ecosystem primed for corridor pilots",
        ],
        size=15,
    )
    _add_textbox(slide, Inches(7.0), Inches(2.05), Inches(5.5), Inches(0.4), "Why Credara", size=14, bold=True, color=BLUE)
    _add_bullets(
        slide,
        Inches(7.0),
        Inches(2.45),
        Inches(5.8),
        Inches(4.0),
        [
            "Live monorepo on Vercel + Amoy — not slideware (test today)",
            "Fail-closed chain: no fake Polygonscan; honest simulated vs anchored",
            "Token-agnostic Smart LC — swap MockUSDC for licensed AED without rebuild",
            "Founder: regulated banking + utilities scale + cryptographic governance (AIRG)",
        ],
        size=15,
    )
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.75), Inches(1.95), Inches(0.05), Inches(4.6))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE7, 0xEA, 0xF1)
    rect.line.fill.background()

    # 5 — How it works
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "HOW IT WORKS", "One trade, five actors — ~3 minutes in Judge Mode")
    steps = [
        ("1", "SME", "Order → invoice → delivery proof"),
        ("2", "Buyer", "Confirm trade obligation"),
        ("3", "SME", "Create receivable · anchor proof on Polygon"),
        ("4", "Financier", "Deal room review · fund Smart LC escrow"),
        ("5", "System", "Verify delivery · release stablecoin to seller"),
    ]
    y = 2.1
    for num, actor, action in steps:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y), Inches(0.5), Inches(0.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BLUE
        rect.line.fill.background()
        _add_textbox(slide, Inches(0.9), Inches(y + 0.06), Inches(0.5), Inches(0.4), num, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(slide, Inches(1.55), Inches(y + 0.05), Inches(2.0), Inches(0.35), actor, size=13, bold=True, color=BLUE)
        _add_textbox(slide, Inches(3.6), Inches(y + 0.02), Inches(8.8), Inches(0.55), action, size=16, color=INK)
        y += 0.92
    example = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.55))
    example.fill.solid()
    example.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    example.line.color.rgb = GREEN
    _add_textbox(
        slide,
        Inches(1.05),
        Inches(6.25),
        Inches(11.2),
        Inches(0.4),
        "Example trade: AED 250k invoice → 80% advance (AED 200k) → Smart LC funded → released same day after verified delivery",
        size=13,
        bold=True,
        color=GREEN,
    )

    # 6 — Traction / proof
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "TRACTION & PROOF", "Built, deployed, and testable — not a concept deck")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.0),
        Inches(5.5),
        Inches(3.5),
        [
            "Production deploy: credara-jet.vercel.app + credara-api.vercel.app",
            "Smart LC factory wired: create → fund → verify → release on Amoy",
            "Auth0 OAuth + password auth · Didit KYB abstraction · Postgres persistence",
            "Slither + Hardhat in CI · role-separated escrow contracts",
            "Judge Mode: guided 3-min demo for reviewers and investors",
        ],
        size=15,
    )
    if WORKSPACE_SHOT.exists():
        slide.shapes.add_picture(str(WORKSPACE_SHOT), Inches(6.6), Inches(2.0), width=Inches(5.9))
        _add_textbox(slide, Inches(6.6), Inches(5.95), Inches(5.9), Inches(0.3), "Enterprise workspace (live product)", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    proof_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.55), Inches(5.5), Inches(1.15))
    proof_box.fill.solid()
    proof_box.fill.fore_color.rgb = NAVY
    proof_box.line.fill.background()
    _add_textbox(slide, Inches(1.05), Inches(5.68), Inches(5.2), Inches(0.35), "Proof moment (demo climax)", size=12, bold=True, color=WHITE)
    if POLYGONSCAN_SHOT.exists():
        slide.shapes.add_picture(str(POLYGONSCAN_SHOT), Inches(1.05), Inches(6.05), width=Inches(2.4))
        _add_textbox(slide, Inches(3.6), Inches(6.05), Inches(2.6), Inches(0.55), "Live Amoy tx: proof anchor + Smart LC release", size=11, color=RGBColor(0xCB, 0xD5, 0xE1))
    else:
        _add_textbox(
            slide,
            Inches(1.05),
            Inches(6.05),
            Inches(5.2),
            Inches(0.55),
            "amoy.polygonscan.com — slot proof-anchor + Smart LC tx screenshot here",
            size=11,
            color=RGBColor(0xCB, 0xD5, 0xE1),
        )

    # 7 — Competition & moat
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "COMPETITION & MOAT", "Infrastructure wedge — not another hash-on-chain demo")
    _add_textbox(slide, Inches(0.85), Inches(2.0), Inches(2.3), Inches(0.35), "Approach", size=11, bold=True, color=BLUE)
    _add_textbox(slide, Inches(3.2), Inches(2.0), Inches(4.5), Inches(0.35), "Limitation", size=11, bold=True, color=MUTED)
    _add_textbox(slide, Inches(7.9), Inches(2.0), Inches(4.5), Inches(0.35), "Credara edge", size=11, bold=True, color=BLUE)
    rows = [
        ("Manual bank LCs", "7–10 days, paper-heavy, SME-unfriendly", "Programmable escrow + verified delivery gate"),
        ("Trade finance fintech", "Off-chain docs; lenders still blind to delivery", "Buyer-confirmed receivable + on-chain anchor"),
        ("Generic RWA / trade chains", "Hash demos; simulated chain links", "Fail-closed product · live Amoy · four-role workspaces"),
        ("Remittance / POS apps", "Consumer wedge; not receivables finance", "B2B settlement slice of UAE corridor"),
    ]
    y = 2.45
    for i, (name, pain, edge) in enumerate(rows):
        _add_comparison_row(slide, y, name, pain, edge, highlight=(i == 2))
        y += 0.78
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(6.35),
        Inches(11.5),
        Inches(0.45),
        "Moat compounds: corridor data + financier network + token-agnostic escrow + evidence governance IP",
        size=13,
        bold=True,
        color=INK,
    )

    # 8 — Business model (simplified)
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "BUSINESS MODEL", "Who pays: financiers and platform users — API expands TAM")
    streams = [
        ("Financing origination", "Fee when receivable is funded via Smart LC marketplace"),
        ("Platform / SaaS", "Verification, deal room, reporting — tiered by volume"),
        ("Settlement / escrow", "Basis points per fund/release vs. multi-day bank LC cost"),
        ("Credit scoring API", "On-chain attestations for external lenders without full UI"),
    ]
    y = 2.2
    for title, desc in streams:
        _add_textbox(slide, Inches(0.9), Inches(y), Inches(3.4), Inches(0.4), title, size=16, bold=True, color=BLUE)
        _add_textbox(slide, Inches(4.4), Inches(y), Inches(8.1), Inches(0.5), desc, size=16, color=INK)
        y += 0.82
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(5.75),
        Inches(11.5),
        Inches(0.55),
        "Land financiers first — verified receivables pull SME supply onto the rail.",
        size=15,
        bold=True,
        color=GREEN,
    )

    # 9 — Unit economics + ask
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "UNIT ECONOMICS & ASK", "Illustrative pilot model — tighten with live pilot quotes")
    _add_stat_card(slide, Inches(0.9), Inches(2.1), Inches(2.5), Inches(1.1), "0.5%", "Platform fee (indicative)")
    _add_stat_card(slide, Inches(3.55), Inches(2.1), Inches(2.5), Inches(1.1), "80%", "Advance rate (indicative)")
    _add_stat_card(slide, Inches(6.2), Inches(2.1), Inches(2.5), Inches(1.1), "AED 1.25k", "Fee on AED 250k invoice")
    _add_stat_card(slide, Inches(8.85), Inches(2.1), Inches(2.8), Inches(1.1), "$15M", "Pilot receivables pipe (Y1 target)")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(3.45),
        Inches(6.2),
        Inches(3.2),
        [
            "Pilot cohort: 3 financiers + 50 UAE SMEs (Jebel Ali supply chains)",
            "Revenue scales with financed volume, not headcount",
            "API-first: third parties integrate rails without Credara UI",
            "Path to $1M ARR: ~$200M annual financed volume at ~50 bps blended take",
        ],
        size=15,
    )
    ask_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.4), Inches(3.35), Inches(5.2), Inches(2.85))
    ask_box.fill.solid()
    ask_box.fill.fore_color.rgb = NAVY
    ask_box.line.fill.background()
    _add_textbox(slide, Inches(7.6), Inches(3.55), Inches(4.8), Inches(0.4), "Raising (illustrative)", size=12, bold=True, color=RGBColor(0x94, 0xA3, 0xB8))
    _add_textbox(slide, Inches(7.6), Inches(3.95), Inches(4.8), Inches(0.55), "$500K pre-seed", size=28, bold=True, color=WHITE)
    _add_bullets(
        slide,
        Inches(7.6),
        Inches(4.65),
        Inches(4.8),
        Inches(2.4),
        [
            "UAE 3-financier pilot + GTM hire",
            "Didit KYB production + Solidity audit",
            "AED stablecoin issuer partnership",
            "18-mo: mainnet + API ecosystem",
        ],
        size=13,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )

    # 10 — GTM
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "GO-TO-MARKET", "One corridor, prove it, then replicate")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(5.8),
        Inches(4.2),
        [
            "Wedge: UAE SME trade-finance slice — Jebel Ali suppliers",
            "Beachhead: 1–2 non-bank financiers + SME cohort",
            "DIFC / challenge ecosystem for intros and pilot design",
            "KPI: 7–10 day LC → same-day proof-to-funding",
            "Expand: new corridors = issuer + IERC20 config, not rebuild",
        ],
        size=16,
    )
    _add_textbox(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.5), "Roadmap", size=14, bold=True, color=BLUE)
    _add_bullets(
        slide,
        Inches(7.2),
        Inches(2.5),
        Inches(5.5),
        Inches(3.8),
        [
            "Now: Amoy demo + Judge Mode (shipped)",
            "Q1–Q2: Financier pilot LOIs + prod KYB",
            "Phase 2: Licensed AED stablecoin partner",
            "Phase 3: Polygon mainnet + partner APIs",
        ],
        size=15,
        color=MUTED,
    )

    # 11 — Settlement roadmap
    slide = prs.slides.add_slide(blank)
    _set_slide_bg(slide, SOFT)
    _add_textbox(slide, Inches(0.9), Inches(0.6), Inches(11), Inches(0.4), "SETTLEMENT ROADMAP", size=11, bold=True, color=BLUE)
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(1.1),
        Inches(11),
        Inches(1.2),
        "Today: MockUSDC on Amoy proves the escrow path.\nTomorrow: regulated AED stablecoin on the same Smart LC rails.",
        size=28,
        bold=True,
        color=INK,
    )
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.65),
        Inches(11),
        Inches(3.8),
        [
            "Smart LC is token-agnostic — production is an issuer swap, not a rebuild",
            "CBUAE PTSR-compliant AED from licensed partner when ready",
            "B2B settlement along UAE corridors — not consumer remittance or POS",
            "Credara is settlement-ready infrastructure; we do not issue AED",
        ],
        size=17,
    )
    _accent_line(slide, top=2.5)
    _add_footer(slide)

    # 12 — Team
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "TEAM", "Founder + banking engineering leadership — regulated delivery at scale")
    _add_textbox(slide, Inches(0.9), Inches(2.05), Inches(5.6), Inches(0.4), "Othniel Obasi — Founder, NOVTIA Ltd (London)", size=14, bold=True, color=BLUE)
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.45),
        Inches(5.6),
        Inches(3.5),
        [
            "Creditville MFBank 40K → 100K+ users",
            "PHED utilities — 5M+ customers, $96M+ revenue growth",
            "AIRG + SURGE v2 → Credara evidence model",
            "Portfolio: NorthBridge, Kairos, ProofSource, Autonix",
        ],
        size=15,
    )
    _add_textbox(slide, Inches(6.9), Inches(2.05), Inches(5.6), Inches(0.4), "Martins Nwanu — Head of IT, Creditville Group", size=14, bold=True, color=BLUE)
    _add_bullets(
        slide,
        Inches(6.9),
        Inches(2.45),
        Inches(5.6),
        Inches(3.5),
        [
            "11-engineer leadership; core banking platform delivery",
            "C Money 30K → 100K+ users; in-house mobile rebuild",
            "Loan ops scaled to ₦800M–₦1B/month disbursements",
            "Credara: production API, workspace UX, deploy discipline",
        ],
        size=15,
    )
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.65), Inches(1.95), Inches(0.05), Inches(4.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE7, 0xEA, 0xF1)
    rect.line.fill.background()
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(5.85),
        Inches(11.5),
        Inches(0.55),
        "Hiring with raise: UAE GTM / trade-finance partnerships · compliance advisor (DIFC)",
        size=14,
        color=MUTED,
    )

    # 13 — Tech appendix (condensed — one slide for diligence)
    slide = prs.slides.add_slide(blank)
    _section_header(slide, "TECH APPENDIX", "For diligence — full architecture in repo README")
    _add_bullets(
        slide,
        Inches(0.9),
        Inches(2.0),
        Inches(5.5),
        Inches(4.5),
        [
            "Monorepo: Next.js · FastAPI · workers · Solidity 0.8.28",
            "ProofRegistry · ReceivableRegistry · SmartLCFactory",
            "CreditScoreAttestation · MockUSDC (IERC20-agnostic)",
            "Outbox relayer · idempotency · AccessControl / Pausable",
        ],
        size=15,
    )
    _add_bullets(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.8),
        Inches(4.5),
        [
            "Non-goals: no remittance MVP · no POS · no AED issuance",
            "Docs: README · SUBMISSION.md · PROCESS_FLOWS.md",
            "API: credara-api.vercel.app/docs",
            "Contracts: Slither CI + Hardhat test suite",
        ],
        size=15,
        color=MUTED,
    )
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.85), Inches(0.05), Inches(4.8))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE7, 0xEA, 0xF1)
    rect.line.fill.background()

    # 14 — Close + ask
    slide = prs.slides.add_slide(blank)
    _set_slide_bg(slide, NAVY)
    _add_textbox(slide, Inches(0.9), Inches(1.85), Inches(11), Inches(0.9), "Get paid in days, not months.", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(2.85),
        Inches(11),
        Inches(0.8),
        "Verified invoices · Polygon proof · Smart LC settlement",
        size=20,
        color=RGBColor(0xCB, 0xD5, 0xE1),
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(3.85),
        Inches(11),
        Inches(0.6),
        "Let's run the 3-minute Judge Mode demo — then talk UAE pilot.",
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(slide, Inches(0.9), Inches(5.0), Inches(11), Inches(0.5), "credara-jet.vercel.app  ·  othnielobasi@gmail.com", size=14, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
    _add_footer(slide, dark=True)

    prs.save(OUT_INVESTOR)
    return OUT_INVESTOR


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate Credara pitch decks")
    parser.add_argument("--investor", action="store_true", help="Build investor edition only")
    parser.add_argument("--judge", action="store_true", help="Build judge/hackathon edition only")
    args = parser.parse_args()
    build_both = not args.investor and not args.judge
    if args.judge or build_both:
        print(f"Wrote {build()}")
    if args.investor or build_both:
        print(f"Wrote {build_investor()}")


if __name__ == "__main__":
    main()
